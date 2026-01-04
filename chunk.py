import os
import nltk
from nltk.tokenize import sent_tokenize, word_tokenize
from collections import defaultdict
from PIL import Image
from pptx import Presentation
from PyPDF2 import PdfReader
from docx import Document
import pytesseract
import pandas as pd

# Ensure nltk resources are downloaded (if not done previously)
nltk.download('punkt')

class UnifiedDataLoader:

    def load(self, file_path: str):
        ext = os.path.splitext(file_path)[1].lower()

        loaders = {
            ".pdf": self._load_pdf,
            ".doc": self._load_docx,
            ".docx": self._load_docx,
            ".ppt": self._load_ppt,
            ".pptx": self._load_ppt,
            ".xls": self._load_excel,
            ".xlsx": self._load_excel,
            ".jpg": self._load_image,
            ".jpeg": self._load_image,
            ".png": self._load_image,
        }

        if ext not in loaders:
            raise ValueError(f"Unsupported file type: {ext}")

        return loaders[ext](file_path)

    def _load_pdf(self, file_path):
        reader = PdfReader(file_path)
        return [
            {
                "text": page.extract_text() or "",
                "metadata": {
                    "source": os.path.basename(file_path),
                    "type": "pdf",
                    "page": i + 1,
                },
            }
            for i, page in enumerate(reader.pages)
        ]

    def _load_docx(self, file_path):
        doc = Document(file_path)
        return [{
            "text": "\n".join(p.text for p in doc.paragraphs),
            "metadata": {
                "source": os.path.basename(file_path),
                "type": "docx",
            },
        }]

    def _load_ppt(self, file_path):
        prs = Presentation(file_path)
        docs = []

        for i, slide in enumerate(prs.slides):
            text = "\n".join(
                shape.text for shape in slide.shapes if hasattr(shape, "text")
            )
            docs.append({
                "text": text,
                "metadata": {
                    "source": os.path.basename(file_path),
                    "type": "pptx",
                    "slide": i + 1,
                },
            })

        return docs

    def _load_excel(self, file_path):
        sheets = pd.read_excel(file_path, sheet_name=None)
        return [
            {
                "text": df.to_string(index=False),
                "metadata": {
                    "source": os.path.basename(file_path),
                    "type": "excel",
                    "sheet": name,
                },
            }
            for name, df in sheets.items()
        ]

    def _load_image(self, file_path):
        image = Image.open(file_path)
        return [{
            "text": pytesseract.image_to_string(image),
            "metadata": {
                "source": os.path.basename(file_path),
                "type": "image",
            },
        }]

    def process_document(self, file_path, chunking_method='sentences'):
        """
        Process document using the specified chunking method.
        chunking_method: 'sentences', 'paragraphs', 'fixed_size'
        """
        # Step 1: Load the content
        docs = self.load(file_path)

        # Step 2: Chunk the content based on the selected method
        processed_docs = []

        for doc in docs:
            text = doc["text"]

            # Choose chunking method
            if chunking_method == 'sentences':
                chunks = self.chunk_by_sentences(text)
            elif chunking_method == 'paragraphs':
                chunks = self.chunk_by_paragraphs(text)
            elif chunking_method == 'fixed_size':
                chunks = self.chunk_by_fixed_size(text)
            else:
                raise ValueError(f"Unknown chunking method: {chunking_method}")

            # Add chunks to the processed_docs list
            for chunk in chunks:
                processed_docs.append({
                    "text": chunk,
                    "metadata": doc["metadata"],
                })

        return processed_docs

    def chunk_by_sentences(self, text, max_chunk_size=500):
        """
        Chunk text into smaller segments based on sentence tokenization.
        """
        sentences = sent_tokenize(text)
        chunks = []
        current_chunk = ""

        for sentence in sentences:
            if len(current_chunk) + len(sentence) > max_chunk_size:
                chunks.append(current_chunk.strip())
                current_chunk = sentence
            else:
                current_chunk += " " + sentence

        if current_chunk:
            chunks.append(current_chunk.strip())

        return chunks

    def chunk_by_paragraphs(self, text, max_chunk_size=500):
        """
        Chunk text by paragraphs. If a paragraph is too large, split it into smaller chunks.
        """
        paragraphs = text.split("\n\n")  # Assuming paragraphs are separated by two newlines
        chunks = []
        current_chunk = ""

        for paragraph in paragraphs:
            if len(current_chunk) + len(paragraph) > max_chunk_size:
                chunks.append(current_chunk.strip())
                current_chunk = paragraph
            else:
                current_chunk += "\n\n" + paragraph

        if current_chunk:
            chunks.append(current_chunk.strip())

        return chunks

    def chunk_by_fixed_size(self, text, chunk_size=500):
        """
        Chunk the text into fixed-size chunks based on the number of characters.
        """
        return [text[i:i + chunk_size] for i in range(0, len(text), chunk_size)]


# Example usage
file_path = "your_document.pdf"  # Replace with the actual path
data_loader = UnifiedDataLoader()

# Choose your chunking method ('sentences', 'paragraphs', 'fixed_size')
processed_data = data_loader.process_document(file_path, chunking_method='paragraphs')

# View processed data (chunks)
for doc in processed_data:
    print(f"Metadata: {doc['metadata']}")
    print(f"Chunk: {doc['text'][:100]}...")  # Show a preview of each chunk
    print()
